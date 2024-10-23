from flask import Flask, request, jsonify
from flask_cors import CORS
import openai
import pyttsx3
import PyPDF2
import  python-docx
from pptx import Presentation
import os

app = Flask(__name__)
CORS(app)

# Set your OpenAI API key here
openai.api_key = 'https://api.openai.com/v1/organization/projects/{project_id}/api_keys/{key_id}'

# Folder to store uploaded files temporarily
UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Initialize pyttsx3 engine for text-to-speech
engine = pyttsx3.init()

# Function to read PDF file content
def extract_pdf_content(file_path):
    with open(file_path, 'rb') as pdf_file:
        reader = PyPDF2.PdfReader(pdf_file)
        text = ''
        for page in reader.pages:
            text += page.extract_text()
        return text

# Function to read DOCX file content
def extract_docx_content(file_path):
    doc = python_docx.Document(file_path)
    text = '\n'.join([para.text for para in doc.paragraphs])
    return text

# Function to read PPTX file content
def extract_pptx_content(file_path):
    prs = Presentation(file_path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

# Handle file upload and parsing
@app.route('/upload-file', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"status": "Error", "message": "No file part"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"status": "Error", "message": "No selected file"}), 400

    # Save the uploaded file
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
    file.save(file_path)

    # Determine file type and extract content
    if file.filename.endswith('.pdf'):
        content = extract_pdf_content(file_path)
    elif file.filename.endswith('.docx'):
        content = extract_docx_content(file_path)
    elif file.filename.endswith('.pptx'):
        content = extract_pptx_content(file_path)
    else:
        return jsonify({"status": "Error", "message": "Unsupported file format"}), 400

    # Optionally, delete the file after processing
    os.remove(file_path)

    return jsonify({"status": "Success", "content": content}), 200

# Handle AI query based on extracted file content
@app.route('/ask-question', methods=['POST'])
def ask_question():
    try:
        data = request.json
        question = data.get('question', '')
        file_content = data.get('content', '')

        if not question or not file_content:
            return jsonify({"status": "Error", "message": "Invalid input"}), 400

        # Call OpenAI API to get a response from ChatGPT
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": f"{question}\nContext: {file_content}"}
            ]
        )

        answer = response['choices'][0]['message']['content']

        # Convert the answer to speech
        engine.say(answer)
        engine.runAndWait()

        return jsonify({"status": "Success", "message": answer}), 200

    except Exception as e:
        return jsonify({"status": "Error", "message": str(e)}), 500

# Vehicle health check via voice
@app.route('/check-vehicle-health', methods=['POST'])
def check_vehicle_health():
    try:
        data = request.json  # Receive JSON data from frontend
        print(f"Received data: {data}")  # Log the data for debugging
        # Here you can process the data and check the vehicle health
        if data:
            # Example health check response
            return jsonify({"status": "Success", "message": "Vehicle is healthy"}), 200
        else:
            return jsonify({"status": "Error", "message": "No data received"}), 400
    except Exception as e:
        return jsonify({"status": "Error", "message": str(e)}), 500

# Test route
@app.route('/hello', methods=['GET'])
def check_vehicle():
    return "Hello from Flask!"

if __name__ == '__main__':
    # Ensure upload folder exists
    if not os.path.exists(UPLOAD_FOLDER):
        os.makedirs(UPLOAD_FOLDER)

    app.run(debug=True)
